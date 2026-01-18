import html
import re
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from html.parser import HTMLParser

class HTMLRunParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.runs = []
        self.current_style = {
            "bold": False, "italic": False, "underline": False, "strike": False,
            "font_size": None, "color_rgb": None, "theme_color": None, "brightness": None
        }
        self.style_stack = []

    def handle_starttag(self, tag, attrs):
        self.style_stack.append(self.current_style.copy())
        attrs_dict = dict(attrs)

        if tag == "br":
            # Treat <br> as data "\x0b" (soft return)
            self.runs.append({
                "text": "\x0b",
                "style": self.current_style.copy()
            })
        elif tag in ["b", "strong"]:
            self.current_style["bold"] = True
        elif tag in ["i", "em"]:
            self.current_style["italic"] = True
        elif tag == "u":
            self.current_style["underline"] = True
        elif tag in ["s", "strike", "del"]:
            self.current_style["strike"] = True
        # Parse new tags: <c v="...">, <sz v="...">
        elif tag == "c":
            val = attrs_dict.get("v", "")
            if val.startswith("#"):
                # RGB
                self.current_style["color_rgb"] = val.replace("#", "").upper()
            elif val.startswith("T"):
                # Theme: T1 or T1:0.5
                parts = val[1:].split(":")
                try:
                    self.current_style["theme_color"] = int(parts[0])
                    if len(parts) > 1:
                        self.current_style["brightness"] = float(parts[1])
                except ValueError:
                    pass

        elif tag == "sz":
            val = attrs_dict.get("v", "")
            try:
                self.current_style["font_size"] = float(val)
            except ValueError:
                pass

        # Legacy parsing support (optional, but good for robustness if mixed)
        elif tag in ["span", "font"]:
             pass # Ignore legacy for now to force new schema usage

    def handle_startendtag(self, tag, attrs):
        # Handle <br /> self-closing
        if tag == "br":
             self.runs.append({
                "text": "\x0b",
                "style": self.current_style.copy()
            })

    def handle_endtag(self, tag):
        if tag == "br":
            return # Ignore </br> if it exists, or handled in starttag
        if self.style_stack:
            self.current_style = self.style_stack.pop()

    def handle_data(self, data):
        if not data:
            return
        self.runs.append({
            "text": data,
            "style": self.current_style.copy()
        })

class PPTXProcessor:
    def __init__(self, filepath, translator):
        self.filepath = filepath
        self.translator = translator
        self.prs = Presentation(filepath)

    def process(self):
        """
        Iterates through all slides, collecting and translating text in batches per slide.
        """
        total_slides = len(self.prs.slides)
        print(f"Processing {total_slides} slides...")

        for i, slide in enumerate(tqdm(self.prs.slides, desc="Translating Slides")):
            self._process_slide(slide, slide_index=i)

    def _process_slide(self, slide, slide_index):
        # Collect tasks for this slide
        tasks = []
        for shape in slide.shapes:
            self._collect_tasks(shape, tasks)

        if not tasks:
            return

        # Separate into Standard and Constrained batches
        standard_tasks = [t for t in tasks if t['context'] == 'standard']
        constrained_tasks = [t for t in tasks if t['context'] == 'constrained']

        if standard_tasks:
            self._process_batch(standard_tasks, self.translator.config.presentation_body_prompt, f"Slide {slide_index+1} (Standard)")

        if constrained_tasks:
            self._process_batch(constrained_tasks, self.translator.config.constrained_text_prompt, f"Slide {slide_index+1} (Constrained)")

    def _process_batch(self, tasks, prompt_template, description):
        batch_items = []

        # Prepare Batch
        for i, task in enumerate(tasks):
            paragraph = task["paragraph"]
            html_text = self._paragraph_to_html(paragraph)

            # Skip empty
            if not html_text.strip():
                continue

            raw_text_length = len("".join([r.text for r in paragraph.runs]))
            max_chars = self._calculate_max_chars(raw_text_length)

            batch_items.append({
                "id": i,
                "text": html_text,
                "limit": max_chars,
                "_task_ref": task # Keep reference to original task
            })

        if not batch_items:
            return

        # Remove _task_ref before sending to LLM
        llm_payload = [{k: v for k, v in item.items() if k != "_task_ref"} for item in batch_items]

        # Translate
        translated_items = self.translator.translate_batch(llm_payload, prompt_template)

        # Validate
        if len(translated_items) != len(batch_items):
            print(f"\n[Error] {description}: Batch count mismatch. Sent {len(batch_items)}, received {len(translated_items)}. Skipping.")
            return

        # Sort by ID to ensure alignment (though LLM should preserve order, IDs are safer)
        # We assume the LLM returns IDs.
        translated_map = {item.get("id"): item.get("translation") for item in translated_items}

        # Apply Translations
        for item in batch_items:
            t_id = item["id"]
            if t_id not in translated_map:
                print(f"\n[Error] {description}: ID {t_id} missing in response. Skipping item.")
                continue

            translated_text = translated_map[t_id]
            paragraph = item["_task_ref"]["paragraph"]
            self._reconstruct_paragraph(paragraph, translated_text)

    def _collect_tasks(self, shape, task_list, context="standard"):
        # Handle Groups (Recursive)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for item in shape.shapes:
                self._collect_tasks(item, task_list, context="constrained")
            return

        # Handle Tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if not cell.text_frame:
                        continue
                    self._collect_text_frame_tasks(cell.text_frame, task_list, context="constrained")
            return

        # Handle Text Frames
        if shape.has_text_frame:
            self._collect_text_frame_tasks(shape.text_frame, task_list, context=context)

    def _collect_text_frame_tasks(self, text_frame, task_list, context="standard"):
        for paragraph in text_frame.paragraphs:
            if not paragraph.text.strip():
                continue

            task_list.append({
                "paragraph": paragraph,
                "context": context
            })

    def _paragraph_to_html(self, paragraph):
        """
        Converts paragraph runs to HTML string.
        """
        parts = []
        for run in paragraph.runs:
            parts.append(self._run_to_html(run))
        return "".join(parts)

    def _run_to_html(self, run):
        # Escape HTML first
        text = html.escape(run.text)
        if not text:
            return ""

        # Convert line breaks/soft returns to <br> to ensure LLM visibility and preservation
        # _x000B_ is the string representation of \x0b in python-pptx text runs sometimes
        text = text.replace("_x000B_", "<br>").replace("\x0b", "<br>").replace("\n", "<br>").replace("\r", "<br>")

        # Color & Size Attributes
        c_tag = None
        sz_tag = None

        # 1. Color (c tag)
        try:
            if run.font.color.type == 1: # RGB
                c_tag = f'<c v="#{run.font.color.rgb}">'
            elif run.font.color.type == 2: # Theme
                # Format: T<id> or T<id>:<brightness>
                t_val = f"T{run.font.color.theme_color}"
                if run.font.color.brightness:
                     t_val += f":{run.font.color.brightness}"
                c_tag = f'<c v="{t_val}">'
        except:
            pass

        # 2. Size (sz tag)
        if run.font.size and run.font.size.pt:
            sz_tag = f'<sz v="{int(run.font.size.pt)}">'

        # Construct tags
        result = text

        # Apply wrappers (Order: Size, Color, Bold, Italic, Underline, Strike)
        # Inner-most should be formatting, outer-most structural?
        # Actually standard HTML nesting doesn't matter too much for parser, but shorter first.

        if run.font.bold:
            result = f"<b>{result}</b>"
        if run.font.italic:
            result = f"<i>{result}</i>"
        try:
            if run.font.underline:
                result = f"<u>{result}</u>"
        except: pass
        try:
            if hasattr(run.font, 'strike') and run.font.strike:
                result = f"<s>{result}</s>"
        except: pass

        if c_tag:
            result = f"{c_tag}{result}</c>"

        if sz_tag:
            result = f"{sz_tag}{result}</sz>"

        return result

    def _reconstruct_paragraph(self, paragraph, html_text):
        if not html_text:
            return

        # Parse HTML
        parser = HTMLRunParser()
        try:
            parser.feed(html_text)
        except Exception as e:
            print(f"HTML Parse Error: {e} | Text: {html_text}")
            # Fallback: Just set text if parse fails
            paragraph.clear()
            paragraph.add_run().text = html.unescape(html_text)
            return

        parsed_runs = parser.runs
        if not parsed_runs:
            return

        # Clear and rebuild
        paragraph.clear()

        for p_run in parsed_runs:
            new_run = paragraph.add_run()
            new_run.text = html.unescape(p_run["text"])
            self._apply_style(new_run, p_run["style"])

    def _apply_style(self, run, style):
        run.font.bold = style["bold"]
        run.font.italic = style["italic"]
        run.font.underline = style["underline"]
        if style["strike"]:
            try:
                if hasattr(run.font, 'strike'):
                    run.font.strike = True
            except: pass

        if style["font_size"]:
            run.font.size = Pt(style["font_size"])

        if style["color_rgb"]:
            try:
                from pptx.dml.color import RGBColor
                run.font.color.rgb = RGBColor.from_string(style["color_rgb"])
            except: pass

        if style["theme_color"]:
            try:
                run.font.color.theme_color = int(style["theme_color"])
                if style["brightness"] is not None:
                    run.font.color.brightness = style["brightness"]
            except: pass

    def _calculate_max_chars(self, original_length):
        ratio = 1.0
        conf = self.translator.config

        s_lang = conf.source_language.lower()
        t_lang = conf.target_language.lower()
        base_ratio = conf.expansion_ratio

        if "japanese" in s_lang and "english" in t_lang:
            ratio = base_ratio
        elif "english" in s_lang and "japanese" in t_lang:
            ratio = 1.0 / base_ratio if base_ratio != 0 else 1.0

        return int(original_length * ratio)

    def save(self, output_path):
        self.prs.save(output_path)
