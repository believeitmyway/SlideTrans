from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from tqdm import tqdm
import math

class LayoutAdjuster:
    def __init__(self, filepath):
        self.filepath = filepath
        self.prs = Presentation(filepath)

    def adjust(self):
        """
        Iterates through slides and adjusts text boxes and tables.
        """
        for slide_idx, slide in enumerate(tqdm(self.prs.slides, desc="Adjusting Layout")):
            self._adjust_slide(slide)

    def save(self, output_path):
        self.prs.save(output_path)

    def _adjust_slide(self, slide):
        # We need to collect all shapes first to do collision detection
        shapes = list(slide.shapes)

        for shape in shapes:
            if shape.has_text_frame:
                self._adjust_text_box(shape, shapes, slide.slide_layout)

            if shape.has_table:
                self._adjust_table(shape)

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._adjust_group(shape)

    def _adjust_text_box(self, shape, all_shapes, layout):
        # Only adjust if it's not a table or chart
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return

        try:
            current_left = shape.left
            current_top = shape.top
            current_width = shape.width
            current_height = shape.height
            current_right = current_left + current_width

            # Find nearest obstruction to the right
            slide_width = self.prs.slide_width
            max_right = slide_width

            # Check other shapes
            for other in all_shapes:
                if other == shape:
                    continue

                # Simple AABB collision check for "in the same horizontal band"
                if other.left >= current_right:
                    if not (other.top + other.height < current_top or other.top > current_top + current_height):
                        if other.left < max_right:
                            max_right = other.left

            # Calculate new width
            # Leave a small margin (e.g., 0.1 inch)
            margin = 91440
            available_width = max_right - current_left - margin

            # 1. Widen the shape if possible
            if available_width > current_width:
                shape.width = int(available_width)

            # Ensure word wrap is on (important for multi-line calculation)
            shape.text_frame.word_wrap = True

            # 2. Manual Font Scaling
            # We must pass the available height too, assuming shape height is fixed/max
            self._apply_manual_fit(shape.text_frame, shape.width, shape.height)

        except Exception as e:
            # print(f"Error adjusting shape: {e}")
            pass

    def _adjust_table(self, shape):
        for row in shape.table.rows:
            row_height = row.height
            for cell in row.cells:
                if cell.text_frame:
                    cell.text_frame.word_wrap = True
                    # Check against cell width and row height
                    self._apply_manual_fit(cell.text_frame, cell.width, row_height)

    def _adjust_group(self, group_shape):
        for shape in group_shape.shapes:
             if shape.has_text_frame:
                 shape.text_frame.word_wrap = True
                 self._apply_manual_fit(shape.text_frame, shape.width, shape.height)

    def _apply_manual_fit(self, text_frame, available_width, available_height):
        """
        Estimates text height based on wrapping and shrinks font if it exceeds available_height.
        """
        if not available_width or available_width <= 0:
            return

        # If height is not provided or zero, we assume it can grow, so we don't shrink?
        # But user said height is fixed.
        if not available_height or available_height <= 0:
            available_height = 99999999 # Treat as infinite if unknown

        # Iterative shrinking? Or calculate once.
        # Calculating exact wrap is hard.
        # Let's do a heuristic:
        # 1. Calculate total length of text in EMUs (if it were one line).
        # 2. Divide by available_width to get estimated lines.
        # 3. Multiply by line_height to get total height.

        total_text_width_linear = self._estimate_total_linear_width(text_frame)

        # Avoid division by zero
        if total_text_width_linear == 0:
            return

        # Average Font Size estimate (weighted? or just max?)
        # We need a representative font size to calculate line height.
        avg_font_size_pt = self._get_max_font_size(text_frame)
        if avg_font_size_pt == 0:
            avg_font_size_pt = 18 # Fallback

        line_height_emu = avg_font_size_pt * 12700 * 1.2 # Approx 1.2 spacing

        # Estimated lines
        # We need to account that words cannot be split easily, so effective width usage is < 100%.
        # Let's assume 90% efficiency.
        effective_width = available_width * 0.95

        estimated_lines = math.ceil(total_text_width_linear / effective_width)

        # If there are explicit paragraphs, each paragraph starts a new line.
        # The linear width sum method underestimates if there are many short paragraphs.
        # Better: Sum linear width per paragraph, calc lines per paragraph.

        estimated_lines = 0
        for paragraph in text_frame.paragraphs:
            p_width = self._estimate_paragraph_linear_width(paragraph)
            if p_width == 0:
                # Empty paragraph = 1 line (blank line)
                lines = 1
            else:
                lines = math.ceil(p_width / effective_width)
            estimated_lines += lines

        estimated_total_height = estimated_lines * line_height_emu

        if estimated_total_height > available_height:
            # We need to shrink.
            # Height is proportional to Font Size (Line Height) * Lines.
            # Lines is proportional to 1 / Font Size (roughly).
            # Wait, LineWidth ~ FontSize.
            # Lines = (TotalChars * FontSize) / BoxWidth.
            # Height = Lines * (FontSize * 1.2)
            # Height = (TotalChars * FontSize / BoxWidth) * (FontSize * 1.2)
            # Height = (TotalChars * 1.2 / BoxWidth) * FontSize^2
            # Height = K * FontSize^2

            # Ratio needed = AvailableHeight / EstimatedHeight
            # (NewSize / OldSize)^2 = Ratio
            # NewSize / OldSize = sqrt(Ratio)

            ratio = available_height / estimated_total_height
            scale_factor = math.sqrt(ratio)

            # Apply safety buffer
            safe_factor = scale_factor * 0.95

            self._scale_font_size(text_frame, safe_factor)

    def _estimate_total_linear_width(self, text_frame):
        total = 0
        for p in text_frame.paragraphs:
            total += self._estimate_paragraph_linear_width(p)
        return total

    def _estimate_paragraph_linear_width(self, paragraph):
        width = 0
        for run in paragraph.runs:
            font_size = run.font.size
            if font_size is None:
                font_size = Pt(18)

            text = run.text
            # Remove newlines for linear calculation?
            # Actually explicit newlines in a run should force line breaks,
            # but usually they are separate paragraphs.
            # If text has \n, it acts like soft break.
            # For simplicity, treat chars as linear flow.

            for char in text:
                if char == '\n':
                    continue # Handle separately?
                is_wide = ord(char) > 255
                factor = 1.0 if is_wide else 0.55
                char_width = font_size.pt * 12700 * factor
                width += char_width
        return width

    def _get_max_font_size(self, text_frame):
        max_size = 0
        for p in text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and r.font.size.pt > max_size:
                    max_size = r.font.size.pt
        return max_size

    def _scale_font_size(self, text_frame, ratio):
        """
        Multiplies the font size of all runs by the ratio.
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    new_size = run.font.size.pt * ratio
                    if new_size < 6:
                        new_size = 6
                    run.font.size = Pt(new_size)
                else:
                    new_size = 18 * ratio
                    if new_size < 6:
                        new_size = 6
                    run.font.size = Pt(new_size)
