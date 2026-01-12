import unittest
from unittest.mock import MagicMock, patch
from src.layout_adjuster import LayoutAdjuster
from pptx.util import Pt

class MockFont:
    def __init__(self, size_pt=12.9):
        self.size = MagicMock()
        self.size.pt = size_pt

class MockRun:
    def __init__(self, text, font_size=12.9):
        self.text = text
        self.font = MockFont(font_size)

class MockParagraph:
    def __init__(self, runs):
        self.runs = runs

class MockTextFrame:
    def __init__(self, paragraphs, width=1000000, height=500000):
        # Width/Height in EMUs.
        # 1 inch = 914400 EMUs.
        # Let's say box is 5 inches wide, 2 inches high.
        self.width = width
        self.height = height # Fixed height
        self.paragraphs = paragraphs
        self.word_wrap = True

class TestFontScaling(unittest.TestCase):
    @patch("src.layout_adjuster.Presentation")
    def test_excessive_shrinking(self, mock_presentation):
        # Scenario from user:
        # 427 English characters.
        # Original Font: 12.9pt.
        # Result was 6pt (Too small).
        # Should be ~11pt.

        text = "A" * 427 # 427 chars
        original_size = 12.9

        run = MockRun(text, original_size)
        paragraph = MockParagraph([run])

        # Let's define a reasonable box width.
        # 427 chars at 11pt (~5.5pt width/char for narrow) -> ~2300pt total linear width.
        # If box is 500pt wide (approx 7 inches), lines = 4.6 lines.
        # Height at 11pt * 1.2 = 13.2pt line height.
        # Total height = 5 * 13.2 = 66pt.

        # Let's set box width/height in EMUs.
        # 1 pt = 12700 EMUs.
        # Width = 500pt = 6,350,000 EMUs.
        # Height = 100pt = 1,270,000 EMUs (Plenty of space for 66pt).

        box_width = 500 * 12700
        box_height = 100 * 12700

        text_frame = MockTextFrame([paragraph], width=box_width, height=box_height)

        # Instantiate LayoutAdjuster (mocking file loading)
        adjuster = LayoutAdjuster("dummy.pptx")

        # Run _apply_manual_fit
        adjuster._apply_manual_fit(text_frame, text_frame.width, text_frame.height)

        new_size = run.font.size.pt
        print(f"Original Size: {original_size}pt")
        print(f"New Size: {new_size}pt")

        # We expect it to FAIL currently (produce tiny font)
        # We want it to be > 10pt.
        self.assertGreater(new_size, 10.0, f"Font shrunk too much: {new_size}")

if __name__ == "__main__":
    unittest.main()
