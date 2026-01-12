from pptx import Presentation
from pptx.util import Inches

def test_cell_width_fix():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes

    rows = 2
    cols = 2
    left = top = Inches(2.0)
    width = height = Inches(2.0)

    table_shape = shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Simulate the logic in layout_adjuster
    print("Iterating table cells...")
    for row in table.rows:
        for col_idx, cell in enumerate(row.cells):
            # The fix: access column width
            try:
                cell_width = table.columns[col_idx].width
                print(f"Cell ({row}, {col_idx}) Width from column: {cell_width}")
            except Exception as e:
                print(f"FAILED: {e}")
                raise

if __name__ == "__main__":
    test_cell_width_fix()
